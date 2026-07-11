import asyncio
from pathlib import Path

import httpx
import fitz

from app.config import settings


async def parse_with_mineru(file_path: str | Path) -> str:
    url = f"{settings.MINERU_URL}/parse"
    file_path = Path(file_path)
    async with httpx.AsyncClient(timeout=300) as client:
        with open(file_path, "rb") as f:
            resp = await client.post(
                url,
                files={"file": (file_path.name, f, _guess_mime(file_path))},
            )
        resp.raise_for_status()
        data = resp.json()
        return data.get("content", "")


def _guess_mime(path: Path) -> str:
    ext = path.suffix.lower()
    return {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".doc": "application/msword",
        ".ppt": "application/vnd.ms-powerpoint",
    }.get(ext, "application/octet-stream")


async def parse_with_fallback(file_path: str | Path) -> str:
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return await asyncio.to_thread(_parse_pdf, file_path)
    elif ext in (".docx", ".doc"):
        return await asyncio.to_thread(_parse_docx, file_path)
    elif ext in (".pptx", ".ppt"):
        return await asyncio.to_thread(_parse_pptx, file_path)
    return ""


def _parse_pdf(path: Path) -> str:
    doc = fitz.open(path)
    parts = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            parts.append(f"--- 第 {page.number + 1} 页 ---\n{text}")
    doc.close()
    return "\n\n".join(parts)


def _parse_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            parts.append("\n".join(texts))
    return "\n\n".join(parts)


async def parse_document(file_path: str | Path) -> str:
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    if settings.MINERU_URL:
        try:
            return await parse_with_mineru(file_path)
        except Exception as e:
            pass

    return await parse_with_fallback(file_path)
